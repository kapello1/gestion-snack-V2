# -*- coding: utf-8 -*-
"""
Génération du PowerPoint TFE — Gestion Snack
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Cm
import copy

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

# ── Palette ───────────────────────────────────────────────────────────────────
NAVY   = RGBColor(0x1A, 0x23, 0x5E)
TEAL   = RGBColor(0x00, 0x7A, 0x8A)
ORANGE = RGBColor(0xE6, 0x7E, 0x22)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
LGRAY  = RGBColor(0xF2, 0xF4, 0xF7)
DGRAY  = RGBColor(0x44, 0x44, 0x44)
MGRAY  = RGBColor(0x88, 0x88, 0x88)
GREEN  = RGBColor(0x27, 0xAE, 0x60)
RED    = RGBColor(0xC0, 0x39, 0x2B)

W = prs.slide_width
H = prs.slide_height

# ── Helpers ──────────────────────────────────────────────────────────────────
def blank_slide():
    layout = prs.slide_layouts[6]  # blank
    return prs.slides.add_slide(layout)

def fill_bg(slide, color: RGBColor):
    from pptx.oxml.ns import qn
    from lxml import etree
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def rect(slide, x, y, w, h, fill_color=None, line_color=None, line_width=Pt(0)):
    from pptx.util import Emu
    shp = slide.shapes.add_shape(1, x, y, w, h)  # MSO_SHAPE_TYPE.RECTANGLE
    shp.line.width = line_width
    if fill_color:
        shp.fill.solid()
        shp.fill.fore_color.rgb = fill_color
    else:
        shp.fill.background()
    if line_color:
        shp.line.color.rgb = line_color
        shp.line.width = line_width if line_width else Pt(1)
    else:
        shp.line.fill.background()
    return shp

def txt(slide, text, x, y, w, h, size=Pt(18), bold=False, color=WHITE,
        align=PP_ALIGN.LEFT, italic=False, wrap=True):
    txb = slide.shapes.add_textbox(x, y, w, h)
    tf  = txb.text_frame
    tf.word_wrap = wrap
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size  = size
    run.font.bold  = bold
    run.font.color.rgb = color
    run.font.italic = italic
    return txb

def bullet_box(slide, items, x, y, w, h, size=Pt(16), color=DGRAY,
               bullet_color=TEAL, title=None, title_color=TEAL):
    txb = slide.shapes.add_textbox(x, y, w, h)
    tf  = txb.text_frame
    tf.word_wrap = True
    first = True
    if title:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        r = p.add_run()
        r.text = title
        r.font.size = size + Pt(2)
        r.font.bold = True
        r.font.color.rgb = title_color
    for item in items:
        p = tf.add_paragraph() if not first else tf.paragraphs[0]
        first = False
        p.space_before = Pt(4)
        r = p.add_run()
        r.text = f"▸  {item}"
        r.font.size  = size
        r.font.color.rgb = color
    return txb

def add_notes(slide, text):
    notes = slide.notes_slide
    tf = notes.notes_text_frame
    tf.text = text

def accent_bar(slide, color=TEAL, height=Inches(0.07)):
    rect(slide, 0, H - height, W, height, fill_color=color)

def slide_number(slide, n, total=15):
    txt(slide, f"{n} / {total}",
        W - Inches(1.2), H - Inches(0.38), Inches(1.0), Inches(0.3),
        size=Pt(11), color=MGRAY, align=PP_ALIGN.RIGHT)

def header_band(slide, title, subtitle=None,
                bg=NAVY, fg=WHITE, sub_fg=LGRAY):
    rect(slide, 0, 0, W, Inches(1.6), fill_color=bg)
    txt(slide, title,
        Inches(0.5), Inches(0.2), W - Inches(1), Inches(0.85),
        size=Pt(28), bold=True, color=fg)
    if subtitle:
        txt(slide, subtitle,
            Inches(0.5), Inches(1.0), W - Inches(1), Inches(0.5),
            size=Pt(15), color=sub_fg, italic=True)

def tag(slide, label, x, y, bg=TEAL, fg=WHITE, size=Pt(13)):
    w = Inches(2.2)
    h = Inches(0.38)
    rect(slide, x, y, w, h, fill_color=bg)
    txt(slide, label, x, y + Pt(2), w, h,
        size=size, bold=True, color=fg, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════════════════════
#  SLIDE 1 — PAGE DE TITRE
# ══════════════════════════════════════════════════════════════════════════════
sl = blank_slide()
fill_bg(sl, NAVY)

# Bande teal décorative
rect(sl, 0, Inches(2.8), W, Inches(0.08), fill_color=TEAL)
rect(sl, 0, Inches(5.5), W, Inches(0.08), fill_color=ORANGE)

# Titre principal
txt(sl, "GESTION SNACK",
    Inches(0.8), Inches(1.0), W - Inches(1.6), Inches(1.1),
    size=Pt(52), bold=True, color=WHITE, align=PP_ALIGN.CENTER)

txt(sl, "Application Web de Gestion d'Établissement de Restauration",
    Inches(0.8), Inches(2.1), W - Inches(1.6), Inches(0.7),
    size=Pt(18), color=TEAL, align=PP_ALIGN.CENTER, italic=True)

# Infos
txt(sl, "Tiegni Gamo Bernard-Joël",
    Inches(0.8), Inches(3.3), W - Inches(1.6), Inches(0.5),
    size=Pt(20), bold=True, color=WHITE, align=PP_ALIGN.CENTER)

txt(sl, "Institut Libre Marie Haps  •  Section Informatique  •  TFE 2025–2026",
    Inches(0.8), Inches(3.9), W - Inches(1.6), Inches(0.5),
    size=Pt(14), color=LGRAY, align=PP_ALIGN.CENTER)

# Stack badges
badges = [
    (Inches(1.8), "React 19"),
    (Inches(3.6), "Spring Boot 3"),
    (Inches(5.7), "PostgreSQL 16"),
    (Inches(7.6), "Groq AI"),
    (Inches(9.4), "ElevenLabs"),
    (Inches(11.1), "Stripe"),
]
for bx, label in badges:
    tag(sl, label, bx, Inches(5.0), bg=TEAL)

txt(sl, "17 juin 2026",
    Inches(0.8), Inches(6.7), W - Inches(1.6), Inches(0.4),
    size=Pt(13), color=MGRAY, align=PP_ALIGN.CENTER)

add_notes(sl,
    "Bonjour, je vais vous présenter mon projet TFE : une application web complète de gestion d'un snack, "
    "développée en React et Spring Boot, déployée dans le cloud. Cette application couvre la gestion des "
    "commandes, du stock, de la caisse, des réservations et intègre un chatbot vocal intelligent.")

# ══════════════════════════════════════════════════════════════════════════════
#  SLIDE 2 — CONTEXTE & PROBLÉMATIQUE
# ══════════════════════════════════════════════════════════════════════════════
sl = blank_slide()
fill_bg(sl, WHITE)
header_band(sl, "Contexte & Problématique", "Pourquoi cette application ?")
accent_bar(sl)
slide_number(sl, 2)

problems = [
    "Commandes perdues ou mal transmises entre salle et cuisine",
    "Aucun suivi du stock en temps réel → ruptures fréquentes",
    "Gestion de caisse manuelle → erreurs de calcul",
    "Aucune donnée analysée pour piloter l'activité",
    "Réservations gérées par téléphone ou carnet papier",
    "Communication salle ↔ cuisine lente et source d'erreurs",
]
bullet_box(sl, problems,
           Inches(0.5), Inches(1.8), Inches(6.2), Inches(5.0),
           size=Pt(15), color=DGRAY, title="Problèmes sans système informatisé")

# Solution box
rect(sl, Inches(7.0), Inches(1.8), Inches(5.8), Inches(5.0),
     fill_color=NAVY)
solutions = [
    "Commandes en temps réel via WebSocket",
    "Stock auto-décrémenté par trigger PostgreSQL",
    "Caisse multi-méthodes (CASH / CARD / Stripe)",
    "Dashboard analytique avec KPIs",
    "Réservations en ligne par les clients",
    "Chatbot vocal IA pour les clients",
]
bullet_box(sl, solutions,
           Inches(7.2), Inches(1.9), Inches(5.4), Inches(4.8),
           size=Pt(15), color=WHITE, title="Notre solution", title_color=TEAL)

add_notes(sl,
    "Un snack sans système informatisé souffre de plusieurs problèmes : les commandes se perdent, "
    "le stock n'est pas suivi, les erreurs de caisse sont fréquentes, et aucune donnée n'est analysée "
    "pour prendre des décisions. Mon application résout ces problèmes avec un système centralisé "
    "accessible depuis n'importe quel navigateur.")

# ══════════════════════════════════════════════════════════════════════════════
#  SLIDE 3 — ÉTUDE DE L'EXISTANT
# ══════════════════════════════════════════════════════════════════════════════
sl = blank_slide()
fill_bg(sl, WHITE)
header_band(sl, "Étude de l'Existant", "Analyse comparative des solutions du marché")
accent_bar(sl)
slide_number(sl, 3)

# 4 cards
cards = [
    ("Square POS",    TEAL,   ["Paiements intégrés", "Interface intuitive"],
                              ["Pas de chatbot IA", "Pas de gestion stock avancée", "Coût élevé"]),
    ("Lightspeed",    ORANGE, ["Gestion complète", "Multi-sites"],
                              ["Trop complexe pour PME", "Inaccessible financièrement", "Pas d'IA"]),
    ("Tiller",        NAVY,   ["Interface tactile", "Rapports"],
                              ["Pas de rôles granulaires", "Pas d'API publique", "Pas de chatbot"]),
    ("Notre projet",  GREEN,  ["Chatbot vocal IA (Groq + ElevenLabs)", "Triggers stock automatiques",
                               "API REST Swagger", "6 rôles granulaires", "Open-source friendly"], []),
]
for i, (name, col, pros, cons) in enumerate(cards):
    cx = Inches(0.3 + i * 3.2)
    rect(sl, cx, Inches(1.7), Inches(3.0), Inches(0.45), fill_color=col)
    txt(sl, name, cx, Inches(1.72), Inches(3.0), Inches(0.4),
        size=Pt(14), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    rect(sl, cx, Inches(2.15), Inches(3.0), Inches(4.9),
         fill_color=LGRAY, line_color=col, line_width=Pt(1.5))
    lines = [f"✓  {p}" for p in pros] + ([" "] if pros and cons else []) + [f"✗  {c}" for c in cons]
    bullet_box(sl, lines, cx + Inches(0.1), Inches(2.2),
               Inches(2.8), Inches(4.7), size=Pt(11.5), color=DGRAY,
               bullet_color=col, title=None)

add_notes(sl,
    "J'ai analysé les solutions existantes. Square POS et Lightspeed sont puissants mais trop coûteux "
    "et trop complexes pour une PME. Aucun n'offre de chatbot vocal IA, de gestion du stock par triggers "
    "automatiques, ou d'API REST documentée. Notre projet se distingue sur ces trois points.")

# ══════════════════════════════════════════════════════════════════════════════
#  SLIDE 4 — ARCHITECTURE TECHNIQUE
# ══════════════════════════════════════════════════════════════════════════════
sl = blank_slide()
fill_bg(sl, WHITE)
header_band(sl, "Architecture Technique", "Architecture N-tier déployée dans le cloud")
accent_bar(sl)
slide_number(sl, 4)

# 3 couches
layers = [
    (Inches(0.4),  TEAL,  "COUCHE PRÉSENTATION",
     "React 19  •  Vite 7  •  Tailwind CSS 3",
     "Vercel — CDN mondial (edge network)"),
    (Inches(4.5),  NAVY,  "COUCHE MÉTIER",
     "Spring Boot 3  •  Java 17  •  WebSocket STOMP",
     "Render — Conteneur Linux (Oregon)"),
    (Inches(8.6),  ORANGE,"COUCHE DONNÉES",
     "PostgreSQL 16  •  ENUMs  •  Triggers",
     "Neon.tech — Serverless (EU Frankfurt)"),
]
for lx, col, title, tech, host in layers:
    rect(sl, lx, Inches(1.8), Inches(3.8), Inches(4.5), fill_color=col)
    txt(sl, title, lx + Inches(0.1), Inches(1.9), Inches(3.6), Inches(0.55),
        size=Pt(12), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    txt(sl, tech,  lx + Inches(0.15), Inches(2.55), Inches(3.5), Inches(0.8),
        size=Pt(13.5), color=WHITE, align=PP_ALIGN.CENTER)
    rect(sl, lx + Inches(0.2), Inches(3.5), Inches(3.4), Inches(0.04),
         fill_color=WHITE)
    txt(sl, "☁  " + host, lx + Inches(0.1), Inches(3.65), Inches(3.6), Inches(0.6),
        size=Pt(12), color=LGRAY, align=PP_ALIGN.CENTER, italic=True)

# Flèches texte
txt(sl, "←  HTTPS  →", Inches(4.0), Inches(3.3), Inches(0.7), Inches(0.4),
    size=Pt(10), color=MGRAY, align=PP_ALIGN.CENTER)
txt(sl, "←  JDBC  →",  Inches(8.1), Inches(3.3), Inches(0.7), Inches(0.4),
    size=Pt(10), color=MGRAY, align=PP_ALIGN.CENTER)

# Services externes
ext_services = ["Groq AI  (chatbot)", "ElevenLabs  (TTS vocal)", "Stripe  (paiements)"]
for i, svc in enumerate(ext_services):
    sx = Inches(0.4 + i * 4.3)
    rect(sl, sx, Inches(6.35), Inches(3.8), Inches(0.7), fill_color=LGRAY)
    txt(sl, svc, sx + Inches(0.1), Inches(6.4), Inches(3.6), Inches(0.55),
        size=Pt(12), color=DGRAY, align=PP_ALIGN.CENTER)

txt(sl, "↑  REST API  ↑", Inches(0.4), Inches(6.1), Inches(3.8), Inches(0.3),
    size=Pt(10), color=MGRAY, align=PP_ALIGN.CENTER)
txt(sl, "↑  REST API  ↑", Inches(4.5), Inches(6.1), Inches(3.8), Inches(0.3),
    size=Pt(10), color=MGRAY, align=PP_ALIGN.CENTER)

add_notes(sl,
    "L'architecture est en 3 couches. Le frontend React est hébergé sur Vercel avec CDN mondial. "
    "Le backend Spring Boot tourne sur Render dans un conteneur Linux. La base PostgreSQL est sur "
    "Neon.tech, serverless, hébergée en Europe pour la conformité RGPD. Les trois déploiements "
    "sont automatiques à chaque push GitHub.")

# ══════════════════════════════════════════════════════════════════════════════
#  SLIDE 5 — TECHNOLOGIES
# ══════════════════════════════════════════════════════════════════════════════
sl = blank_slide()
fill_bg(sl, WHITE)
header_band(sl, "Choix Technologiques", "Justification de chaque technologie retenue")
accent_bar(sl)
slide_number(sl, 5)

techs = [
    ("Spring Boot 3.5 + Java 17", NAVY,
     "Standard industriel Java, auto-configuration, écosystème Spring Data JPA + WebSocket"),
    ("React 19 + Vite 7",         TEAL,
     "DOM virtuel, hooks, HMR ultra-rapide, CDN Vercel — SPA moderne et maintenable"),
    ("PostgreSQL 16 (Neon.tech)",  ORANGE,
     "ENUMs natifs, triggers avancés, serverless EU (RGPD), scaling automatique"),
    ("Groq API — LLaMA 3.1-8b",   NAVY,
     "Latence < 1s (LPU), 5x plus rapide qu'OpenAI — essentiel pour le chatbot vocal"),
    ("ElevenLabs TTS",             TEAL,
     "Voix française naturelle, modèle turbo v2.5, API REST simple, fallback navigateur prévu"),
    ("Stripe",                     ORANGE,
     "Standard industrie, 3D Secure natif, webhooks, données bancaires jamais sur notre serveur"),
]
for i, (name, col, desc) in enumerate(techs):
    row = i // 2
    col_idx = i % 2
    tx = Inches(0.4 + col_idx * 6.4)
    ty = Inches(1.9 + row * 1.65)
    rect(sl, tx, ty, Inches(5.9), Inches(1.5), fill_color=LGRAY,
         line_color=col, line_width=Pt(2))
    rect(sl, tx, ty, Inches(0.18), Inches(1.5), fill_color=col)
    txt(sl, name, tx + Inches(0.28), ty + Inches(0.1), Inches(5.5), Inches(0.45),
        size=Pt(14), bold=True, color=NAVY)
    txt(sl, desc, tx + Inches(0.28), ty + Inches(0.55), Inches(5.5), Inches(0.85),
        size=Pt(12), color=DGRAY, wrap=True)

add_notes(sl,
    "Groq plutôt qu'OpenAI car la latence est 5x plus faible — essentiel pour le chatbot vocal. "
    "PostgreSQL pour les ENUMs natifs et les triggers avancés. BCrypt pour les mots de passe, "
    "sessions côté serveur — pas de JWT nécessaire pour ce périmètre.")

# ══════════════════════════════════════════════════════════════════════════════
#  SLIDE 6 — LES 6 RÔLES
# ══════════════════════════════════════════════════════════════════════════════
sl = blank_slide()
fill_bg(sl, WHITE)
header_band(sl, "Les 6 Rôles du Système", "Contrôle d'accès par rôle (RBAC) — 26 user stories")
accent_bar(sl)
slide_number(sl, 6)

roles = [
    ("👤  ADMIN",        NAVY,   "Dashboard analytique, gestion personnel,\nfournisseurs, catalogue, audit log"),
    ("💰  CAISSIER",     TEAL,   "Encaissement (CASH/CARD/Stripe),\nticket PDF, annulation commande"),
    ("🍽  SERVEUR",      ORANGE, "Prise de commande en salle,\ngestion tables & réservations"),
    ("👨‍🍳  CUISINIER",    NAVY,   "Réception commandes temps réel,\nstatuts préparation, alertes stock"),
    ("📱  CLIENT",       TEAL,   "Menu en ligne, commandes, réservations,\npaiement Stripe, chatbot vocal, avis"),
    ("🚚  FOURNISSEUR",  ORANGE, "Commandes d'approvisionnement,\nconfirmation livraisons, produits en alerte"),
]
for i, (name, col, desc) in enumerate(roles):
    row = i // 3
    ci  = i % 3
    rx  = Inches(0.4 + ci * 4.3)
    ry  = Inches(1.85 + row * 2.5)
    rect(sl, rx, ry, Inches(4.0), Inches(2.2), fill_color=col)
    txt(sl, name, rx + Inches(0.1), ry + Inches(0.15), Inches(3.8), Inches(0.55),
        size=Pt(15), bold=True, color=WHITE)
    txt(sl, desc, rx + Inches(0.15), ry + Inches(0.75), Inches(3.7), Inches(1.3),
        size=Pt(12.5), color=LGRAY, wrap=True)

add_notes(sl,
    "Le système distingue 6 rôles : Admin, Caissier, Serveur, Cuisinier, Client et Fournisseur. "
    "Chaque rôle a une interface dédiée et un périmètre fonctionnel précis, géré par un système "
    "RBAC côté backend Spring Boot.")

# ══════════════════════════════════════════════════════════════════════════════
#  SLIDE 7 — TITRE DÉMO
# ══════════════════════════════════════════════════════════════════════════════
sl = blank_slide()
fill_bg(sl, NAVY)
rect(sl, 0, Inches(2.5), W, Inches(0.1), fill_color=TEAL)
rect(sl, 0, Inches(4.9), W, Inches(0.1), fill_color=ORANGE)

txt(sl, "DÉMONSTRATION", Inches(0.5), Inches(1.0), W - Inches(1), Inches(1.2),
    size=Pt(58), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
txt(sl, "Application en production — Live Demo",
    Inches(0.5), Inches(2.7), W - Inches(1), Inches(0.6),
    size=Pt(20), color=TEAL, align=PP_ALIGN.CENTER, italic=True)

scenarios = [
    "1. Dashboard Admin & Tableau de bord",
    "2. Commande en temps réel (WebSocket)",
    "3. Paiement Stripe & Ticket PDF",
    "4. Chatbot Vocal (Groq + ElevenLabs)",
    "5. Réservation client en ligne",
]
for i, sc in enumerate(scenarios):
    txt(sl, sc,
        Inches(2.5), Inches(3.4 + i * 0.65), Inches(9.0), Inches(0.6),
        size=Pt(16), color=LGRAY, align=PP_ALIGN.CENTER)

slide_number(sl, 7)
add_notes(sl,
    "DÉMO LIVE — Ouvrez l'application dans le navigateur. "
    "Ayez deux onglets : un connecté Serveur/Admin, un connecté Cuisinier pour montrer WebSocket.")

# ══════════════════════════════════════════════════════════════════════════════
#  SLIDE 8 — DÉMO : DASHBOARD ADMIN
# ══════════════════════════════════════════════════════════════════════════════
sl = blank_slide()
fill_bg(sl, WHITE)
header_band(sl, "Démo 1 — Dashboard Admin",
            "KPIs, gestion catalogue, alertes stock en temps réel")
accent_bar(sl)
slide_number(sl, 8)

items_left = [
    "Connexion Admin → tableau de bord principal",
    "Chiffre d'affaires : jour / semaine / mois",
    "Graphiques des tendances de ventes",
    "Top 5 produits les plus vendus",
    "Journal d'audit : toutes les actions tracées",
]
bullet_box(sl, items_left,
           Inches(0.5), Inches(1.8), Inches(6.0), Inches(4.5),
           size=Pt(15), color=DGRAY, title="Ce qu'on montre", title_color=NAVY)

# Encadré Alertes Stock
rect(sl, Inches(7.0), Inches(1.8), Inches(5.8), Inches(4.5), fill_color=LGRAY,
     line_color=ORANGE, line_width=Pt(2))
txt(sl, "⚡  Alertes Stock Automatiques",
    Inches(7.2), Inches(1.9), Inches(5.4), Inches(0.55),
    size=Pt(14), bold=True, color=ORANGE)
al_items = [
    "Trigger PostgreSQL déclenché à chaque commande",
    "Si stock < seuil → INSERT dans stock_alerts",
    "Notification WebSocket instantanée vers Admin",
    "Badge rouge sur le dashboard",
    "Aucun code applicatif — garanti par la BD",
]
bullet_box(sl, al_items,
           Inches(7.2), Inches(2.55), Inches(5.4), Inches(3.6),
           size=Pt(13), color=DGRAY)

add_notes(sl,
    "MONTREZ : connectez-vous en Admin → dashboard → catalogue produits → alertes stock. "
    "Soulignez que les alertes sont générées par un trigger PostgreSQL sans intervention applicative.")

# ══════════════════════════════════════════════════════════════════════════════
#  SLIDE 9 — DÉMO : COMMANDE TEMPS RÉEL
# ══════════════════════════════════════════════════════════════════════════════
sl = blank_slide()
fill_bg(sl, WHITE)
header_band(sl, "Démo 2 — Commande en Temps Réel",
            "Communication WebSocket STOMP entre Serveur et Cuisinier")
accent_bar(sl)
slide_number(sl, 9)

# Flux
steps = [
    (TEAL,   "1  SERVEUR",       "Sélectionne table 3\nAjoute 2x Burger + 1x Frites"),
    (NAVY,   "2  BACKEND",       "POST /api/orders\nINSERT orders + order_items\ntrigger_update_stock"),
    (ORANGE, "3  WEBSOCKET",     "/topic/orders\nMessage diffusé au rôle COOK"),
    (GREEN,  "4  CUISINIER",     "Notification reçue\n< 100ms après la validation"),
]
for i, (col, title, detail) in enumerate(steps):
    sx = Inches(0.4 + i * 3.2)
    rect(sl, sx, Inches(1.85), Inches(2.9), Inches(4.6), fill_color=col)
    txt(sl, title, sx + Inches(0.1), Inches(1.95), Inches(2.7), Inches(0.55),
        size=Pt(14), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    txt(sl, detail, sx + Inches(0.1), Inches(2.6), Inches(2.7), Inches(3.5),
        size=Pt(13), color=LGRAY, align=PP_ALIGN.CENTER, wrap=True)
    if i < 3:
        txt(sl, "→", Inches(3.15 + i * 3.2), Inches(3.4), Inches(0.3), Inches(0.5),
            size=Pt(28), bold=True, color=MGRAY, align=PP_ALIGN.CENTER)

rect(sl, Inches(0.4), Inches(6.55), W - Inches(0.8), Inches(0.6),
     fill_color=LGRAY)
txt(sl, "💡  Ouvrez 2 onglets : Serveur + Cuisinier — la notification arrive en moins de 100ms",
    Inches(0.5), Inches(6.58), W - Inches(1.0), Inches(0.5),
    size=Pt(13), bold=True, color=NAVY, align=PP_ALIGN.CENTER)

add_notes(sl,
    "MONTREZ : ouvrez deux onglets. Dans l'onglet Serveur créez une commande. "
    "L'onglet Cuisinier reçoit la notification WebSocket instantanément. "
    "Changez le statut de la commande depuis la cuisine.")

# ══════════════════════════════════════════════════════════════════════════════
#  SLIDE 10 — DÉMO : PAIEMENT STRIPE
# ══════════════════════════════════════════════════════════════════════════════
sl = blank_slide()
fill_bg(sl, WHITE)
header_band(sl, "Démo 3 — Paiement Stripe & Ticket PDF",
            "Encaissement sécurisé + justificatif PDF généré côté client")
accent_bar(sl)
slide_number(sl, 10)

stripe_steps = [
    ("Caissier", TEAL,   "Sélectionne la commande\nchoisit méthode STRIPE"),
    ("Frontend", NAVY,   "Crée un PaymentIntent\nvia /api/payments/stripe"),
    ("Stripe",   ORANGE, "stripe.confirmCardPayment()\nCarte test : 4242 4242 4242 4242"),
    ("Backend",  GREEN,  "status = CLOSED\nINSERT transaction\njsPDF → ticket PDF"),
]
for i, (actor, col, desc) in enumerate(stripe_steps):
    sx = Inches(0.4 + i * 3.2)
    rect(sl, sx, Inches(1.85), Inches(2.9), Inches(3.5), fill_color=col)
    txt(sl, actor, sx + Inches(0.1), Inches(1.95), Inches(2.7), Inches(0.5),
        size=Pt(14), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    txt(sl, desc,  sx + Inches(0.1), Inches(2.55), Inches(2.7), Inches(2.6),
        size=Pt(13), color=LGRAY, align=PP_ALIGN.CENTER, wrap=True)
    if i < 3:
        txt(sl, "→", Inches(3.15 + i * 3.2), Inches(3.2), Inches(0.3), Inches(0.5),
            size=Pt(28), bold=True, color=MGRAY, align=PP_ALIGN.CENTER)

rect(sl, Inches(0.4), Inches(5.55), Inches(6.1), Inches(1.6), fill_color=LGRAY)
txt(sl, "🔒  Sécurité Stripe",
    Inches(0.6), Inches(5.65), Inches(5.7), Inches(0.45),
    size=Pt(13), bold=True, color=NAVY)
txt(sl, "Les données bancaires ne transitent JAMAIS par notre serveur.\n"
        "Stripe tokenise directement — conformité PCI DSS garantie.",
    Inches(0.6), Inches(6.15), Inches(5.7), Inches(0.85),
    size=Pt(12), color=DGRAY, wrap=True)

rect(sl, Inches(6.8), Inches(5.55), Inches(6.1), Inches(1.6), fill_color=LGRAY)
txt(sl, "📄  Ticket PDF (jsPDF)",
    Inches(7.0), Inches(5.65), Inches(5.7), Inches(0.45),
    size=Pt(13), bold=True, color=NAVY)
txt(sl, "Généré entièrement côté navigateur React.\n"
        "Contient : date, articles, total, méthode de paiement.",
    Inches(7.0), Inches(6.15), Inches(5.7), Inches(0.85),
    size=Pt(12), color=DGRAY, wrap=True)

add_notes(sl,
    "MONTREZ : Caissier → sélectionner commande → STRIPE → taper 4242 4242 4242 4242 → "
    "confirmer → ticket PDF généré et téléchargé automatiquement.")

# ══════════════════════════════════════════════════════════════════════════════
#  SLIDE 11 — DÉMO : CHATBOT VOCAL
# ══════════════════════════════════════════════════════════════════════════════
sl = blank_slide()
fill_bg(sl, NAVY)
header_band(sl, "Démo 4 — Chatbot Vocal Intelligent",
            "Web Speech API  →  Groq LLaMA 3.1  →  ElevenLabs TTS", bg=NAVY)
accent_bar(sl, color=ORANGE)
slide_number(sl, 11)

# Machine à états
states = [
    (MGRAY,  "IDLE",       "Chatbot en attente"),
    (TEAL,   "LISTENING",  "Micro actif\nWeb Speech API\nlang: fr-FR"),
    (ORANGE, "PROCESSING", "Groq API\nLLaMA 3.1-8b-instant\nlatence < 1s"),
    (GREEN,  "SPEAKING",   "ElevenLabs TTS\neleven_turbo_v2_5\nVoix Adam (fr)"),
]
for i, (col, state, desc) in enumerate(states):
    sx = Inches(0.5 + i * 3.15)
    rect(sl, sx, Inches(1.9), Inches(2.8), Inches(3.8), fill_color=col)
    txt(sl, state, sx + Inches(0.1), Inches(2.0), Inches(2.6), Inches(0.55),
        size=Pt(16), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    txt(sl, desc, sx + Inches(0.1), Inches(2.65), Inches(2.6), Inches(2.8),
        size=Pt(13), color=LGRAY, align=PP_ALIGN.CENTER, wrap=True)
    if i < 3:
        txt(sl, "→", Inches(3.1 + i * 3.15), Inches(3.4), Inches(0.3), Inches(0.5),
            size=Pt(28), bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# Flèche retour
txt(sl, "← Retour automatique en fin d'audio (événement onended)",
    Inches(0.5), Inches(5.85), W - Inches(1.0), Inches(0.4),
    size=Pt(12), color=TEAL, align=PP_ALIGN.CENTER, italic=True)

# Infos bas
infos = [
    "🛡  Anti-écho : micro désactivé pendant la lecture TTS",
    "⚡  Turbo v2.5 : latence divisée par 2 vs modèle standard",
    "🔄  Fallback : speechSynthesis navigateur si ElevenLabs indisponible",
]
for i, info in enumerate(infos):
    txt(sl, info, Inches(0.5 + i * 4.3), Inches(6.35), Inches(4.1), Inches(0.7),
        size=Pt(12), color=WHITE, align=PP_ALIGN.CENTER, wrap=True)

add_notes(sl,
    "MONTREZ : espace client → activer le chatbot → parler au micro → "
    "posez 'Qu'est-ce que vous avez au menu ?'. "
    "Montrez la transcription, la réponse texte générée par Groq, puis la lecture vocale ElevenLabs.")

# ══════════════════════════════════════════════════════════════════════════════
#  SLIDE 12 — POINTS TECHNIQUES CLÉS
# ══════════════════════════════════════════════════════════════════════════════
sl = blank_slide()
fill_bg(sl, WHITE)
header_band(sl, "Points Techniques Clés", "Triggers PostgreSQL  •  WebSocket  •  Sécurité BCrypt")
accent_bar(sl)
slide_number(sl, 12)

blocks = [
    (NAVY,  "Triggers PostgreSQL",
     ["AFTER INSERT on order_items",
      "→ UPDATE products SET stock = stock - qty",
      "Si stock < alert_threshold :",
      "→ INSERT INTO stock_alerts",
      "Cohérence garantie par la BD elle-même"]),
    (TEAL,  "WebSocket STOMP",
     ["Broker Spring Boot intégré",
      "Topics par rôle : /topic/orders, /topic/alerts",
      "SockJS fallback pour compatibilité",
      "Clients React s'abonnent au démarrage",
      "Latence < 100ms en conditions normales"]),
    (ORANGE,"Sécurité",
     ["BCrypt facteur 10 — aucun mot de passe en clair",
      "RBAC sur chaque endpoint Spring Boot",
      "HTTPS partout (certificats SSL Render/Vercel)",
      "Clés API en variables d'environnement",
      "Requêtes paramétrées — anti SQL injection"]),
]
for i, (col, title, items) in enumerate(blocks):
    bx = Inches(0.4 + i * 4.3)
    rect(sl, bx, Inches(1.85), Inches(4.0), Inches(5.0), fill_color=LGRAY,
         line_color=col, line_width=Pt(2.5))
    rect(sl, bx, Inches(1.85), Inches(4.0), Inches(0.55), fill_color=col)
    txt(sl, title, bx + Inches(0.1), Inches(1.88), Inches(3.8), Inches(0.5),
        size=Pt(13), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    bullet_box(sl, items, bx + Inches(0.15), Inches(2.5),
               Inches(3.7), Inches(4.2), size=Pt(12.5), color=DGRAY)

add_notes(sl,
    "Triggers PostgreSQL : quand un article est commandé, un trigger décrémente automatiquement "
    "le stock. Si le stock passe sous le seuil, un second trigger génère une alerte. "
    "Tout cela sans code applicatif — c'est la base de données qui garantit la cohérence. "
    "BCrypt pour les mots de passe, RBAC sur tous les endpoints, HTTPS partout.")

# ══════════════════════════════════════════════════════════════════════════════
#  SLIDE 13 — GITHUB & CI/CD
# ══════════════════════════════════════════════════════════════════════════════
sl = blank_slide()
fill_bg(sl, WHITE)
header_band(sl, "GitHub & CI/CD", "Intégration continue — déploiement automatique en < 3 minutes")
accent_bar(sl)
slide_number(sl, 13)

# Pipeline
pipeline = [
    (NAVY,   "git push\nmain",    "Développeur"),
    (TEAL,   "GitHub\nWebhook",   "Notification"),
    (ORANGE, "Render\nBuild",     "mvn package\n→ JAR déployé"),
    (GREEN,  "Vercel\nBuild",     "npm run build\n→ CDN distribué"),
]
for i, (col, title, desc) in enumerate(pipeline):
    px = Inches(0.5 + i * 3.1)
    rect(sl, px, Inches(2.0), Inches(2.7), Inches(1.0), fill_color=col)
    txt(sl, title, px + Inches(0.1), Inches(2.05), Inches(2.5), Inches(0.9),
        size=Pt(14), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    rect(sl, px, Inches(3.1), Inches(2.7), Inches(2.1), fill_color=LGRAY)
    txt(sl, desc, px + Inches(0.1), Inches(3.2), Inches(2.5), Inches(1.9),
        size=Pt(13), color=DGRAY, align=PP_ALIGN.CENTER, wrap=True)
    if i < 3:
        txt(sl, "→", Inches(3.02 + i * 3.1), Inches(2.3), Inches(0.35), Inches(0.5),
            size=Pt(24), bold=True, color=MGRAY)

txt(sl, "⏱  Production en ligne en moins de 3 minutes après chaque push",
    Inches(0.5), Inches(5.35), W - Inches(1.0), Inches(0.5),
    size=Pt(14), bold=True, color=TEAL, align=PP_ALIGN.CENTER)

rect(sl, Inches(0.5), Inches(5.9), W - Inches(1.0), Inches(1.2), fill_color=LGRAY)
txt(sl, "Variables d'environnement sensibles gérées dans les dashboards Render/Vercel — "
        "jamais dans le code source. Fichier .env exclu via .gitignore.",
    Inches(0.7), Inches(6.0), W - Inches(1.4), Inches(1.0),
    size=Pt(13), color=DGRAY, align=PP_ALIGN.CENTER, wrap=True)

add_notes(sl,
    "Tout le code est sur GitHub avec un historique de commits régulier. "
    "Render et Vercel sont connectés au repo : chaque push sur la branche main "
    "déclenche un redéploiement automatique en moins de 3 minutes. "
    "C'est une intégration continue réelle, pas simulée — exigence TFE satisfaite.")

# ══════════════════════════════════════════════════════════════════════════════
#  SLIDE 14 — DIFFICULTÉS RENCONTRÉES
# ══════════════════════════════════════════════════════════════════════════════
sl = blank_slide()
fill_bg(sl, WHITE)
header_band(sl, "Difficultés Rencontrées", "Problèmes réels et solutions apportées")
accent_bar(sl)
slide_number(sl, 14)

challenges = [
    (RED,    "🎤  Anti-écho Chatbot Vocal",
     "Problème : ElevenLabs lisait une réponse → le micro captait l'audio "
     "et relançait une nouvelle requête Groq en boucle.",
     "Solution : désactiver le microphone pendant la lecture TTS. "
     "Timer de sécurité 15s (silence) + timer global 30s (blocage)."),
    (ORANGE, "⚡  Latence TTS trop élevée",
     "Problème : le modèle ElevenLabs standard ajoutait 2-3 secondes "
     "d'attente — expérience vocale non naturelle.",
     "Solution : migration vers eleven_turbo_v2_5 → latence divisée "
     "par 2. Découpage en phrases courtes pour un streaming perçu."),
    (NAVY,   "📊  Diagramme MPD trop grand (400 HTTP)",
     "Problème : le MPD avec 14 tables générait une URL encodée "
     "PlantUML de 8000+ caractères → header HTTP trop large.",
     "Solution : découper le MPD en 2 diagrammes de 7 tables chacun."),
]
for i, (col, title, prob, sol) in enumerate(challenges):
    cy = Inches(1.85 + i * 1.7)
    rect(sl, Inches(0.4), cy, W - Inches(0.8), Inches(1.55), fill_color=LGRAY,
         line_color=col, line_width=Pt(2))
    rect(sl, Inches(0.4), cy, Inches(0.18), Inches(1.55), fill_color=col)
    txt(sl, title, Inches(0.7), cy + Inches(0.08), Inches(4.0), Inches(0.45),
        size=Pt(13), bold=True, color=col)
    txt(sl, f"Problème : {prob}", Inches(0.7), cy + Inches(0.55),
        Inches(5.9), Inches(0.9), size=Pt(11.5), color=DGRAY, wrap=True)
    txt(sl, f"Solution : {sol}", Inches(6.9), cy + Inches(0.08),
        Inches(6.0), Inches(1.35), size=Pt(11.5), color=NAVY, wrap=True)

add_notes(sl,
    "Présentez ces 3 difficultés comme des preuves de votre capacité à résoudre des problèmes réels. "
    "Le jury apprécie que vous ayez rencontré et surmonté des obstacles techniques concrets.")

# ══════════════════════════════════════════════════════════════════════════════
#  SLIDE 15 — CONCLUSION & PERSPECTIVES
# ══════════════════════════════════════════════════════════════════════════════
sl = blank_slide()
fill_bg(sl, NAVY)
rect(sl, 0, Inches(1.5), W, Inches(0.07), fill_color=TEAL)
accent_bar(sl, color=ORANGE)

txt(sl, "Conclusion & Perspectives",
    Inches(0.5), Inches(0.2), W - Inches(1), Inches(1.1),
    size=Pt(30), bold=True, color=WHITE)

# Réalisé
rect(sl, Inches(0.4), Inches(1.7), Inches(5.9), Inches(4.6), fill_color=TEAL)
txt(sl, "✅  Ce qui a été réalisé",
    Inches(0.6), Inches(1.8), Inches(5.5), Inches(0.5),
    size=Pt(14), bold=True, color=WHITE)
done = [
    "Application complète déployée en production",
    "6 rôles  •  26 fonctionnalités  •  14 tables",
    "API REST documentée via Swagger UI",
    "Notifications temps réel (WebSocket STOMP)",
    "Chatbot vocal IA (Groq + ElevenLabs)",
    "Paiements Stripe  •  Export PDF  •  Triggers BD",
    "CI/CD automatique GitHub → Render + Vercel",
]
bullet_box(sl, done, Inches(0.6), Inches(2.4), Inches(5.5), Inches(3.7),
           size=Pt(13), color=WHITE)

# Perspectives
rect(sl, Inches(7.0), Inches(1.7), Inches(5.9), Inches(4.6), fill_color=ORANGE)
txt(sl, "🚀  Améliorations futures",
    Inches(7.2), Inches(1.8), Inches(5.5), Inches(0.5),
    size=Pt(14), bold=True, color=WHITE)
future = [
    "Application mobile React Native",
    "Authentification 2FA (TOTP)",
    "Système de fidélité client",
    "Confirmation réservation par SMS/email",
    "Dashboard prédictif (Machine Learning)",
    "Support multi-établissements",
    "Mode hors-ligne (PWA)",
]
bullet_box(sl, future, Inches(7.2), Inches(2.4), Inches(5.5), Inches(3.7),
           size=Pt(13), color=WHITE)

txt(sl, "Merci de votre attention — Questions ?",
    Inches(0.5), Inches(6.55), W - Inches(1.0), Inches(0.6),
    size=Pt(18), bold=True, color=TEAL, align=PP_ALIGN.CENTER)

slide_number(sl, 15)
add_notes(sl,
    "En conclusion, l'application est complète et déployée en production. "
    "Elle couvre les 6 rôles métier avec 26 fonctionnalités, une API REST documentée via Swagger, "
    "des notifications temps réel, et un chatbot vocal unique. "
    "Pour les évolutions futures : application mobile, 2FA, système de fidélité, "
    "dashboard prédictif par machine learning.")

# ── Sauvegarde ─────────────────────────────────────────────────────────────
output = r"d:\Projet TFE\gestion-snack\fichiers\Presentation_TFE_Gestion_Snack.pptx"
prs.save(output)
print(f"[OK] Présentation générée : {output}")
print(f"     15 slides  •  Notes de présentation incluses")
